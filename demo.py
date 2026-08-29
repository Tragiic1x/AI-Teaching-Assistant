import os
from groq import Groq
import tiktoken
from PyPDF2 import PdfReader
from pptx import Presentation
import streamlit as st

api_key = os.getenv("GROQ_API_KEY")
client = Groq(api_key=api_key)

MODEL = "openai/gpt-oss-120b"
TEMPERATURE = 0.4
MAX_TOKENS = 1500
TOKEN_BUDGET = 6000

SYSTEM_PROMPT = (
    "You are a helpful study assistant. You are given a student's notes or "
    "lecture slides. You either (1) answer questions strictly using only the "
    "provided notes content, saying clearly if the answer isn't in the notes, "
    "or (2) generate quiz questions based on the notes to help the student "
    "study. Be clear, concise, and educational."
)

def shorten_text(text, max_tokens =4000):
    tokens = ENCODING.encode(text)
    if len(tokens) > max_tokens:
        tokens = tokens[:max_tokens]
        return ENCODING.decode(tokens)
    return text

def load_text(file, filename):

    end = filename.lower().rsplit(".", 1)[-1]

    if end == "pdf":
        reader = PdfReader(file)
        text = ""

        for page in reader.pages:
            text += page.extract_text()

        return text 

    elif end == "pptx":
        prs  = Presentation(file)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + " "
                    text += "\n"

        return text

    elif end == "txt":
        return file.read().decode("utf-8")
    
    else:
        return None

def get_encoding(model):

    try:
        return tiktoken.encoding_for_model(model)
    except KeyError:
        print(f"Warning: Tokenizer for model '{model}' not found. Falling back to 'cl100k_base'.")
        return tiktoken.get_encoding("cl100k_base")

ENCODING = get_encoding(MODEL)

def count_tokens(text):
    return len(ENCODING.encode(text))

def total_tokens_used(messages):

    try:
        return sum(count_tokens(msg["content"]) for msg in messages)
    except Exception as e:
        print(f"[token count error]: {e}")
        return 0

def enforce_token_budget(messages, budget=TOKEN_BUDGET):

    try:
        while total_tokens_used(messages) > budget:
            if len(messages) <= 2:
                break
            messages.pop(1)
    except Exception as e:
        print(f"[token budget error]; {e}")

def chat(user_input, messages):
    messages.append({"role": "user", "content": user_input})

    try:

        response = client.chat.completions.create(
            model=MODEL,
            messages=messages,
            temperature=TEMPERATURE,
            max_tokens=MAX_TOKENS
        )
        reply = response.choices[0].message.content
    except Exception as e:
        messages.pop()

        return f"⚠️ Something went wrong talking to the AI: {e}"

    messages.append({"role": "assistant", "content": reply})
    enforce_token_budget(messages)
    return reply

st.set_page_config(page_title="AI Teaching Assistant", page_icon="📚")
st.title("📚 AI Teaching Assistant")
st.caption("Upload your notes or slides, then ask questions or generate a quiz.")

if "messages" not in st.session_state:
    st.session_state.messages = None

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

uploaded_file = st.file_uploader("Upload your notes (PDF, PPTX, or TXT)", type=["pdf", "pptx", "txt"])

if uploaded_file is not None and st.session_state.messages is None:

    notes_text = load_text(uploaded_file, uploaded_file.name)

    if notes_text is None:
        st.error("Unsupported file type. Please upload a PDF, PPTX, or TXT file.")
    elif not notes_text.strip():
        st.error("Couldn't find any readable text in that file.")
    else:
        notes_text = shorten_text(notes_text, max_tokens=4000)
        system_prompt = SYSTEM_PROMPT + f"\n\n--- STUDENT NOTES ---\n{notes_text}\n--- END NOTES ---"
        st.session_state.messages = [{"role": "system", "content": system_prompt}]
        st.success(f"Loaded {len(notes_text.split())} words from {uploaded_file.name}")

if st.session_state.messages is not None:

    st.divider()

    num_of_q = st.number_input("Number of quiz questions", min_value=1, max_value=15, value=5)

    if st.button("🎯 Generate Quiz"):
        prompt = (
            f"Generate {num_of_q} quiz questions (mix of multiple choice and short "
            f"answer) based strictly on the notes provided. Include the correct "
            f"answer clearly labeled at the end of each question."
        )
        reply = chat(prompt, st.session_state.messages)
        st.session_state.chat_history.append(("assistant", reply))

    st.divider()

    for role, text in st.session_state.chat_history:
        with st.chat_message(role):
            st.markdown(text)

    user_input = st.chat_input("What would you like to know?")

    if user_input:
        st.session_state.chat_history.append(("user", user_input))
        with st.chat_message("user"):
            st.markdown(user_input)

        reply = chat (user_input, st.session_state.messages)

        with st.chat_message("assistant"):
            st.markdown(reply)
        st.session_state.chat_history.append(("assistant", reply))

else:
    st.info("👆 Upload a file to get started.")
